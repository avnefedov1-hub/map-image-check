"""Generate Map Image Check architecture presentation (PPTX)."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUTPUT = Path(__file__).resolve().parent / "MapImageCheck_Architecture.pptx"

ACCENT = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x47, 0x55, 0x69)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG = RGBColor(0xF1, 0xF5, 0xF9)


def _set_run(run, *, size: int, bold: bool = False, color: RGBColor = DARK) -> None:
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"


def _add_title_bar(slide, title: str, subtitle: str | None = None) -> None:
    bar = slide.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(13.333), Inches(1.05)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.18), Inches(12.2), Inches(0.55))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    _set_run(r, size=28, bold=True, color=WHITE)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.55), Inches(0.62), Inches(12.2), Inches(0.35))
        stf = sub.text_frame
        stf.clear()
        sp = stf.paragraphs[0]
        sr = sp.add_run()
        sr.text = subtitle
        _set_run(sr, size=14, color=RGBColor(0xDB, 0xEA, 0xFE))


def _add_bullets(slide, items: list[str], *, top: float = 1.35, left: float = 0.65) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(12.0), Inches(5.8))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for index, item in enumerate(items):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = item
        _set_run(r, size=18, color=DARK)


def _add_code_block(slide, text: str, *, top: float = 1.35) -> None:
    shape = slide.shapes.add_shape(1, Inches(0.55), Inches(top), Inches(12.2), Inches(5.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LIGHT_BG
    shape.line.color.rgb = RGBColor(0xD7, 0xDE, 0xE8)

    box = slide.shapes.add_textbox(Inches(0.75), Inches(top + 0.15), Inches(11.8), Inches(5.2))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for index, line in enumerate(text.strip().splitlines()):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = line
        _set_run(r, size=14, color=DARK)
        r.font.name = "Consolas"


def _title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK
    bg.line.fill.background()

    accent = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.18))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.5), Inches(1.2))
    tf = title.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = "Map Image Check"
    _set_run(r, size=44, bold=True, color=WHITE)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.5), Inches(1.0))
    stf = sub.text_frame
    sp = stf.paragraphs[0]
    sr = sp.add_run()
    sr.text = "Архитектура и приложение"
    _set_run(sr, size=26, color=RGBColor(0x93, 0xC5, 0xFD))

    meta = slide.shapes.add_textbox(Inches(0.8), Inches(5.6), Inches(11.5), Inches(0.8))
    mtf = meta.text_frame
    mp = mtf.paragraphs[0]
    mr = mp.add_run()
    mr.text = "Версия 1.1.0  •  av_nefedov  •  Python / Tkinter / OpenCV / SQLite / Ollama"
    _set_run(mr, size=14, color=MUTED)


def _content_slide(prs: Presentation, title: str, subtitle: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, title, subtitle)
    _add_bullets(slide, bullets)


def _diagram_slide(prs: Presentation, title: str, subtitle: str, diagram: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(slide, title, subtitle)
    _add_code_block(slide, diagram)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    _title_slide(prs)

    _content_slide(
        prs,
        "Назначение",
        "Зачем нужно приложение",
        [
            "Автоматический поиск изображений топографических и terrain-карт на дисках и в сети.",
            "Снижение ручного просмотра тысяч файлов за счёт гибридной классификации.",
            "Централизованное хранение найденных карт в SQLite с предпросмотром и метаданными.",
            "Поддержка локальных дисков, папок, компьютеров Active Directory (UNC C$, D$).",
            "Дообучение ML-модели по пользовательским меткам «карта / не карта».",
        ],
    )

    _content_slide(
        prs,
        "Ключевые возможности",
        "Что умеет Map Image Check",
        [
            "Сканирование с пропуском уже известных путей в базе данных.",
            "Гибридный пайплайн: эвристика OpenCV → ML → LLM (серая зона).",
            "Просмотр базы с фильтрами по компьютеру, ML-оценке и LLM-вердикту.",
            "Подробный LLM-анализ на русском через Ollama (одиночный и пакетный).",
            "Экспорт CSV, горячие клавиши, светлая/тёмная тема интерфейса.",
        ],
    )

    _diagram_slide(
        prs,
        "Архитектура (верхний уровень)",
        "Слои приложения",
        """
┌─────────────────────────────────────────────────────────────┐
│  GUI (Tkinter)          gui_app.py  +  search_settings_dialog│
└───────────────────────────────┬─────────────────────────────┘
                                │
        ┌───────────────────────┼───────────────────────┐
        ▼                       ▼                       ▼
┌───────────────┐     ┌─────────────────┐     ┌─────────────────┐
│ scan_drives   │     │ hybrid_pipeline │     │  llm_analysis   │
│ ad_remote     │────▶│ + detector      │────▶│  (Ollama API)   │
└───────────────┘     │ + ml_classifier │     └────────┬────────┘
                      └────────┬────────┘              │
                               ▼                         │
                      ┌─────────────────┐              │
                      │   image_store   │◀─────────────┘
                      │   (SQLite ORM)  │
                      └────────┬────────┘
                               ▼
                      map_image_check.sqlite3
        """,
    )

    _content_slide(
        prs,
        "Модули",
        "Ответственность компонентов",
        [
            "gui_app — главное окно, потоки сканирования/LLM, очередь событий GUI.",
            "scan_drives — обход файловой системы, фильтры размера, CLI-сканер.",
            "ad_remote — список компьютеров AD, проверка online, сбор UNC-путей.",
            "detector — извлечение 10 признаков OpenCV, эвристический score.",
            "ml_classifier — LogisticRegression + StandardScaler, joblib-модель.",
            "hybrid_pipeline — оркестрация этапов и HybridDecision.",
            "llm_analysis — Ollama /api/chat, yes/no и подробный анализ.",
            "image_store — CRUD, индекс путей, метки пользователя, фильтры БД.",
        ],
    )

    _diagram_slide(
        prs,
        "Поток данных",
        "От сканирования до базы",
        """
Пользователь → Выбор источников (диски / папки / AD)
      │
      ▼
_walk_images → фильтр размера → уже в БД? → пропуск
      │
      ▼
classify_image_path (эвристика → ML → LLM?)
      │
      ├── не карта → следующий файл
      │
      └── карта → save_detected_image
                    ├── images (JPEG BLOB)
                    └── heuristic_results (score, features)
      │
      ▼ (опционально)
LLM серая зона / пакетный анализ → llm_results
      │
      ▼
GUI: список карт, предпросмотр, фильтры, метки
        """,
    )

    _diagram_slide(
        prs,
        "Гибридный пайплайн",
        "classify_image_path() — hybrid_pipeline.py",
        """
Эвристика OpenCV (10 признаков, score 0..1)
      │
      ├─ score < T_low (0.30) ──────────────▶ ОТКЛОНИТЬ
      │
      ▼ ML-модель обучена?
      ├─ ml_score ≥ T_accept (0.65) ───────▶ ПРИНЯТЬ (карта)
      ├─ ml_score ≤ T_reject (0.35) ───────▶ ОТКЛОНИТЬ
      │
      ▼ Серая зона (между порогами)
      ├─ LLM включён → yes/no Ollama ──────▶ вердикт LLM
      └─ LLM выключен ─────────────────────▶ ОТКЛОНИТЬ

Fallback без ML: score ≥ 0.45 (MAP_SCORE_THRESHOLD)
        """,
    )

    _content_slide(
        prs,
        "Эвристика и ML",
        "detector.py + ml_classifier.py",
        [
            "Признаки: плотность контуров/рёбер, энтропия цвета, прямые линии, «земляные» тона и др.",
            "Версия детектора: heuristic-v2; порог по умолчанию — 0.45.",
            "ML: StandardScaler + LogisticRegression, сохранение в map_classifier.joblib.",
            "Обучение: минимум 8 образцов на класс; hold-out accuracy при ≥ 20 метках.",
            "Метки пользователя хранятся в user_labels и используются для переобучения.",
        ],
    )

    _content_slide(
        prs,
        "LLM-анализ",
        "llm_analysis.py + Ollama",
        [
            "Модель по умолчанию: qwen2.5vl:7b (vision); endpoint: http://127.0.0.1:11434.",
            "Быстрый yes/no — серая зона при скане (статус gray_zone в БД).",
            "Подробный анализ на русском — вкладка «LLM-анализ» (статус completed).",
            "Пакетный режим: «Проанализировать все в БД» с прогресс-баром.",
            "Поле is_topographic_map — отдельный вердикт «топографическая карта: да/нет».",
            "Фильтры в базе: LLM: карта / не карта / не проверено.",
        ],
    )

    _diagram_slide(
        prs,
        "Схема базы данных",
        "SQLite — map_image_check.sqlite3",
        """
images                          heuristic_results
├── id (PK)                     ├── image_id (PK, FK)
├── source_path, source_host    ├── is_map, score, threshold
├── scan_scope                  ├── detector_version
├── image_bytes (JPEG BLOB)     └── summary_features_json
├── sha256 (UNIQUE)                    (features, ml_score, decision_source)
└── width, height

llm_results                     user_labels
├── image_id (PK, FK)           ├── sha256 (UNIQUE)
├── status (completed/          ├── label (0/1)
│   gray_zone/failed)           └── features_json
├── model_name, analysis_text
├── is_topographic_map
└── structured_json

Лимит BLOB: 50 MiB; JPEG, long side ≤ 2048 px
        """,
    )

    _content_slide(
        prs,
        "Удалённое сканирование",
        "ad_remote.py — Active Directory и UNC",
        [
            "Загрузка списка компьютеров из AD (PowerShell / LDAP).",
            "Проверка доступности: SMB (порт 445) и ping, до 24 параллельных проверок.",
            "Формирование UNC-путей: \\\\COMPUTER\\C$, \\\\COMPUTER\\D$ и др.",
            "Различение серверов и рабочих станций (AdComputerRecord.is_server).",
            "Пользователь выбирает компьютеры и шары в настройках сканирования.",
            "source_host в БД сохраняет имя удалённого компьютера для фильтрации.",
        ],
    )

    _content_slide(
        prs,
        "GUI и сценарии",
        "MapScannerGui — gui_app.py",
        [
            "Сканирование: выбор режима → фоновый поток → список найденных карт.",
            "База данных (Ctrl+B): фильтры по хосту, ML %, LLM-вердикту.",
            "Детали: предпросмотр, классификация, метки, вкладка LLM-анализ.",
            "Настройки (Ctrl+): пороги T_low/T_accept/T_reject, LLM, тема.",
            "Delete — удаление записи из БД (файл на диске не трогается).",
        ],
    )

    _content_slide(
        prs,
        "Технологический стек",
        "Зависимости и платформа",
        [
            "Python 3.10+, Windows (UNC, AD, GUI).",
            "Tkinter / ttk — интерфейс; threading + queue — фоновые задачи.",
            "OpenCV (headless) + NumPy — обработка изображений.",
            "scikit-learn + joblib — ML-модель.",
            "SQLite (stdlib) — персистентность.",
            "Ollama REST API — локальный LLM с поддержкой vision.",
            "PowerShell / LDAP — интеграция с Active Directory.",
        ],
    )

    _content_slide(
        prs,
        "Развёртывание",
        "Запуск, тесты, сборка",
        [
            "pip install -r map_image_check/requirements.txt",
            "python -m map_image_check.gui_app",
            "Тесты: python -m unittest discover -s map_image_check -p test_*.py",
            "Артефакты: map_image_check.sqlite3, map_classifier.joblib (gitignore).",
            "Сборка: PyInstaller (MapImageCheck.spec), Inno Setup installer.",
            "Репозиторий: github.com/avnefedov1-hub/map-image-check",
        ],
    )

    closing = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title_bar(closing, "Итог", "Map Image Check")
    box = closing.shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(11.8), Inches(3.5))
    tf = box.text_frame
    tf.word_wrap = True
    lines = [
        "Многоуровневая классификация снижает нагрузку на LLM и ускоряет скан.",
        "SQLite даёт единый каталог карт с фильтрами и повторным использованием.",
        "ML + пользовательские метки адаптируют систему под ваши данные.",
        "Ollama добавляет объяснимый анализ на русском для спорных случаев.",
    ]
    for index, line in enumerate(lines):
        p = tf.paragraphs[0] if index == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(16)
        r = p.add_run()
        r.text = line
        _set_run(r, size=22, color=DARK)

    prs.save(OUTPUT)
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Created: {path}")
